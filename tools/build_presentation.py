from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Color palette ─────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0F, 0x17, 0x2A)   # slide background
MID_BLUE    = RGBColor(0x1E, 0x3A, 0x8A)   # header bar
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)   # accent / highlight
GREEN       = RGBColor(0x05, 0x96, 0x69)   # success / positive
ORANGE      = RGBColor(0xF5, 0x9E, 0x0B)   # warning
RED         = RGBColor(0xDC, 0x26, 0x26)   # danger
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF1, 0xF5, 0xF9)
YELLOW      = RGBColor(0xFE, 0xF9, 0xC3)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=DARK_BLUE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def header_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.3, MID_BLUE)
    txt(slide, title, 0.4, 0.12, 12, 0.65, size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.75, 12, 0.45, size=14, color=ACCENT_BLUE, italic=True)

def slide_num(slide, n):
    txt(slide, str(n), 12.7, 7.1, 0.5, 0.3, size=11, color=RGBColor(0x64,0x74,0x8B), align=PP_ALIGN.RIGHT)

def bullet_box(slide, items, l, t, w, h, size=15, color=WHITE, indent=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        prefix = '  • ' if indent else ''
        run.text = prefix + item
        run.font.size  = Pt(size)
        run.font.color.rgb = color
    return txb

def metric_box(slide, label, value, unit, l, t, w=2.8, h=1.4,
               bg_color=MID_BLUE, val_color=ACCENT_BLUE):
    rect(slide, l, t, w, h, bg_color)
    txt(slide, value, l+0.1, t+0.1, w-0.2, 0.7, size=32, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txt(slide, unit,  l+0.1, t+0.72, w-0.2, 0.3, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, label, l+0.1, t+1.0,  w-0.2, 0.35, size=12, bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 7.5, DARK_BLUE)
rect(s, 0, 2.8, 13.33, 0.06, ACCENT_BLUE)
txt(s, '🏥', 5.9, 0.5, 1.5, 1.2, size=54, align=PP_ALIGN.CENTER)
txt(s, 'AI-Powered Diabetes Prediction System', 0.5, 1.7, 12.3, 1.0,
    size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, 'Early Detection · Digital Healthcare · Machine Learning',
    0.5, 2.65, 12.3, 0.5, size=16, color=ACCENT_BLUE, align=PP_ALIGN.CENTER, italic=True)
txt(s, 'A Full-Stack Healthcare Platform for Diabetes Risk Screening',
    0.5, 3.2, 12.3, 0.5, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
rect(s, 3.5, 4.1, 6.33, 0.06, MID_BLUE)
txt(s, 'Flask · Python · Gradient Boosting ML · SQLite/PostgreSQL · Chapa Payments',
    0.5, 4.3, 12.3, 0.4, size=12, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.CENTER)
txt(s, 'GitHub: github.com/tasheayansa6/diabetes-prediction-system',
    0.5, 4.85, 12.3, 0.4, size=12, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
txt(s, 'May 2026', 0.5, 6.9, 12.3, 0.4, size=12,
    color=RGBColor(0x64,0x74,0x8B), align=PP_ALIGN.CENTER)
slide_num(s, 1)
print('Slide 1 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Agenda', 'What we will cover today')
items = [
    '1.  Problem Statement — Why this project?',
    '2.  Real-World Healthcare Challenges',
    '3.  Our Solution Overview',
    '4.  System Architecture',
    '5.  Machine Learning Model',
    '6.  Role-Based Clinical Workflow',
    '7.  Payment Integration',
    '8.  Security & Compliance',
    '9.  Testing & Validation',
    '10. Live Demo & Results',
    '11. Impact & Conclusion',
]
bullet_box(s, items, 0.6, 1.5, 12, 5.5, size=16)
slide_num(s, 2)
print('Slide 2 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM: DIABETES STATISTICS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'The Problem: Diabetes is a Global Crisis', 'Why early detection matters')
metric_box(s, 'Adults with Diabetes', '537M', 'worldwide (IDF 2021)', 0.4, 1.5, 2.9, 1.5, MID_BLUE, ACCENT_BLUE)
metric_box(s, 'Undiagnosed Cases', '~50%', 'never detected early', 3.5, 1.5, 2.9, 1.5, RGBColor(0x7C,0x3A,0xED), WHITE)
metric_box(s, 'Projected by 2045', '783M', 'if trend continues', 6.6, 1.5, 2.9, 1.5, RGBColor(0xDC,0x26,0x26), WHITE)
metric_box(s, 'Annual Deaths', '6.7M', 'diabetes-related', 9.7, 1.5, 2.9, 1.5, ORANGE, DARK_BLUE)
txt(s, 'In Ethiopia and developing countries:', 0.6, 3.3, 12, 0.4, size=16, bold=True, color=ACCENT_BLUE)
bullet_box(s, [
    'Most patients diagnosed only when complications are severe (kidney failure, blindness, amputation)',
    'No systematic early screening in primary care clinics',
    'Paper-based records — no data-driven decision support for doctors',
    'Preventable deaths and enormous cost to the healthcare system',
], 0.6, 3.75, 12, 2.8, size=14)
slide_num(s, 3)
print('Slide 3 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PROBLEM: PAPER-BASED WORKFLOW
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Problem 2: Paper-Based, Fragmented Workflow', 'The traditional hospital process')
rect(s, 0.4, 1.5, 5.8, 5.5, MID_BLUE)
txt(s, '❌  BEFORE (Traditional)', 0.6, 1.55, 5.4, 0.45, size=15, bold=True, color=RED)
bullet_box(s, [
    'Nurse writes vitals on paper forms',
    'Lab results recorded in notebooks',
    'Doctor makes decisions without full history',
    'Handwritten prescriptions — illegible',
    'No audit trail for patient actions',
    'Each department works in isolation',
    'Duplicate tests — previous results unavailable',
    'Long waiting times searching for records',
], 0.6, 2.05, 5.4, 4.7, size=13, color=LIGHT_GRAY)
rect(s, 6.8, 1.5, 6.1, 5.5, RGBColor(0x05,0x46,0x2A))
txt(s, '✅  AFTER (This System)', 7.0, 1.55, 5.7, 0.45, size=15, bold=True, color=GREEN)
bullet_box(s, [
    'Nurse records vitals digitally in 2 minutes',
    'Lab results auto-fill patient health form',
    'Doctor sees complete history in one screen',
    'Digital prescriptions — pharmacist verifies online',
    'Every action logged with timestamp and IP',
    'Real-time notifications between departments',
    'No duplicate tests — full history available',
    'Same-day diagnosis and treatment',
], 7.0, 2.05, 5.7, 4.7, size=13, color=LIGHT_GRAY)
slide_num(s, 4)
print('Slide 4 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PROBLEM: PATIENT JOURNEY BEFORE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'The Old Patient Journey', '5 days, multiple trips, high chance of error')
steps = [
    ('Day 1', 'Patient arrives', 'Waits 2 hours\nNurse writes vitals on paper'),
    ('Day 2', 'Doctor visit', 'Orders blood test\nPaper request sent to lab'),
    ('Day 3', 'Lab result', 'Written in notebook\nPatient collects paper'),
    ('Day 4', 'Return to doctor', 'Doctor reads paper\nMakes diagnosis'),
    ('Day 5', 'Prescription', 'Handwritten Rx\nPharmacist cannot read it'),
]
colors = [MID_BLUE, RGBColor(0x1D,0x4E,0xD8), RGBColor(0x7C,0x3A,0xED),
          RGBColor(0xDC,0x26,0x26), ORANGE]
for i, (day, title, detail) in enumerate(steps):
    x = 0.4 + i * 2.55
    rect(s, x, 1.5, 2.3, 1.0, colors[i])
    txt(s, day, x+0.05, 1.55, 2.2, 0.35, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x+0.05, 1.9, 2.2, 0.5, size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(s, '→', x+2.3, 1.85, 0.25, 0.4, size=20, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    rect(s, x, 2.7, 2.3, 2.0, RGBColor(0x1E,0x29,0x3B))
    txt(s, detail, x+0.1, 2.75, 2.1, 1.9, size=12, color=LIGHT_GRAY)
rect(s, 0.4, 5.0, 12.5, 0.7, RGBColor(0x7F,0x1D,0x1D))
txt(s, '⚠  TOTAL: 5 days · Multiple trips · Paper forms · High error rate · No audit trail',
    0.6, 5.1, 12.1, 0.5, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
slide_num(s, 5)
print('Slide 5 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SOLUTION OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Our Solution: AI-Powered Digital Healthcare Platform', 'One system. All roles. Real-time.')
boxes = [
    ('🤖', 'ML Prediction', 'Gradient Boosting\n88.79% Accuracy\n97.06% ROC-AUC', ACCENT_BLUE),
    ('👥', 'Role-Based\nWorkflow', '6 roles: Patient\nDoctor, Nurse, Lab\nPharmacist, Admin', MID_BLUE),
    ('🔔', 'Real-Time\nAlerts', 'WebSocket push\nHigh-risk alerts\nInstant notifications', RGBColor(0x7C,0x3A,0xED)),
    ('💳', 'Payment\nIntegration', 'Chapa · Cash\nInsurance\nDigital receipts', GREEN),
    ('🔒', 'HIPAA\nSecurity', 'JWT Auth · Encryption\nAudit logs\nRate limiting', RGBColor(0xDC,0x26,0x26)),
    ('📊', 'Reports &\nAnalytics', 'PDF reports\nSystem stats\nModel governance', ORANGE),
]
for i, (icon, title, detail, color) in enumerate(boxes):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.3
    y = 1.6 + row * 2.6
    rect(s, x, y, 3.9, 2.3, color)
    txt(s, icon + '  ' + title, x+0.15, y+0.1, 3.6, 0.7, size=16, bold=True, color=WHITE)
    txt(s, detail, x+0.15, y+0.75, 3.6, 1.4, size=12, color=LIGHT_GRAY)
slide_num(s, 6)
print('Slide 6 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — NEW PATIENT JOURNEY
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'The New Patient Journey', 'Same day · Zero paper · Full audit trail')
steps2 = [
    ('Register', 'Patient registers\nNurse notified\nAdded to queue'),
    ('Vitals', 'Nurse records\nBP, BMI, Age\nDigitally in 2 min'),
    ('Lab Test', 'Doctor orders\nLab enters result\nAuto-fills form'),
    ('Payment', 'Patient pays\nChapa/Cash\nInstant confirm'),
    ('ML Predict', 'AI analyses\n8 health features\nRisk in seconds'),
    ('Result', 'Risk level shown\nDoctor alerted\nPrescription sent'),
]
colors2 = [MID_BLUE, RGBColor(0x05,0x96,0x69), RGBColor(0x7C,0x3A,0xED),
           ORANGE, ACCENT_BLUE, GREEN]
for i, (title, detail) in enumerate(steps2):
    x = 0.3 + i * 2.15
    rect(s, x, 1.5, 1.9, 0.7, colors2[i])
    txt(s, title, x+0.05, 1.58, 1.8, 0.5, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 5:
        txt(s, '→', x+1.9, 1.7, 0.25, 0.4, size=18, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    rect(s, x, 2.4, 1.9, 2.2, RGBColor(0x1E,0x29,0x3B))
    txt(s, detail, x+0.1, 2.5, 1.7, 2.0, size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
rect(s, 0.3, 4.85, 12.7, 0.7, RGBColor(0x05,0x46,0x2A))
txt(s, '✅  TOTAL: Same day · Zero paper · Full audit trail · Real-time alerts · AI-powered',
    0.5, 4.95, 12.3, 0.5, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
slide_num(s, 7)
print('Slide 7 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'System Architecture', '3-tier web application with ML integration')
layers = [
    (0.4, 1.5, 12.5, 1.1, ACCENT_BLUE,
     'FRONTEND  (HTML · CSS · JavaScript · Tailwind)',
     'Patient Portal | Doctor Portal | Nurse | Lab | Pharmacist | Admin'),
    (0.4, 2.8, 12.5, 1.1, MID_BLUE,
     'BACKEND  (Flask / Python · REST API · JWT Auth · WebSocket)',
     'Auth Routes | Patient | Doctor | Nurse | Lab | Payment | Admin Routes'),
    (0.4, 4.1, 5.9, 1.1, RGBColor(0x7C,0x3A,0xED),
     'ML SERVICE  (Gradient Boosting)',
     'Accuracy 89.25% | ROC-AUC 97.06%'),
    (6.5, 4.1, 6.4, 1.1, RGBColor(0x05,0x46,0x2A),
     'DATABASE  (SQLite / PostgreSQL)',
     'Users | Patients | Vitals | Predictions | Payments'),
]
for (l, t, w, h, color, title, sub) in layers:
    rect(s, l, t, w, h, color)
    txt(s, title, l+0.15, t+0.08, w-0.3, 0.45, size=14, bold=True, color=WHITE)
    txt(s, sub,   l+0.15, t+0.52, w-0.3, 0.5,  size=11, color=LIGHT_GRAY)
for y in [2.65, 3.95]:
    txt(s, '▼', 6.4, y, 0.5, 0.2, size=14, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
txt(s, '▼', 3.1, 3.95, 0.5, 0.2, size=14, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
txt(s, '▼', 9.5, 3.95, 0.5, 0.2, size=14, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
slide_num(s, 8)
print('Slide 8 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Technology Stack', 'Production-grade tools and frameworks')
techs = [
    ('Backend', ['Python 3.11', 'Flask 3.x', 'SQLAlchemy ORM', 'Flask-SocketIO', 'Flask-Mail', 'JWT Authentication']),
    ('Frontend', ['HTML5 / CSS3', 'JavaScript ES6+', 'Tailwind CSS', 'Bootstrap Icons', 'Chart.js 4.4', 'WebSocket client']),
    ('ML / Data', ['scikit-learn 1.8', 'Gradient Boosting', 'pandas / numpy', 'joblib serialization', 'StandardScaler', '1,068 training samples']),
    ('Database', ['SQLite (dev)', 'PostgreSQL (prod)', 'Flask-Migrate', 'Alembic migrations', 'WAL mode', 'Connection pooling']),
    ('Payments', ['Chapa API', 'TeleBirr / CBE Birr', 'Cash / Insurance', 'Webhook verification', 'PDF receipts', 'Invoice generation']),
    ('Security', ['bcrypt hashing', 'Fernet encryption', 'HIPAA audit logs', 'Rate limiting', 'CSP headers', 'HSTS / HTTPS']),
]
for i, (cat, items) in enumerate(techs):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.3
    y = 1.5 + row * 2.7
    rect(s, x, y, 3.9, 0.45, MID_BLUE)
    txt(s, cat, x+0.1, y+0.05, 3.7, 0.35, size=14, bold=True, color=WHITE)
    bullet_box(s, items, x+0.1, y+0.5, 3.7, 2.1, size=12, color=LIGHT_GRAY)
slide_num(s, 9)
print('Slide 9 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ML MODEL: DATASET
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Machine Learning: Dataset', 'Pima Indians + Frankfurt Diabetes Dataset')
metric_box(s, 'Total Samples', '1,068', 'training records', 0.4, 1.5, 2.8, 1.4)
metric_box(s, 'Non-Diabetic', '650', '60.9% of dataset', 3.4, 1.5, 2.8, 1.4, RGBColor(0x05,0x46,0x2A), GREEN)
metric_box(s, 'Diabetic', '418', '39.1% of dataset', 6.4, 1.5, 2.8, 1.4, RGBColor(0x7F,0x1D,0x1D), RED)
metric_box(s, 'Features', '8', 'clinical variables', 9.4, 1.5, 2.8, 1.4, RGBColor(0x4C,0x1D,0x95), WHITE)
txt(s, '8 Input Features:', 0.6, 3.2, 12, 0.4, size=15, bold=True, color=ACCENT_BLUE)
features = [
    'Pregnancies — number of times pregnant',
    'Glucose — plasma glucose concentration (mg/dL)',
    'BloodPressure — diastolic blood pressure (mmHg)',
    'SkinThickness — triceps skin fold thickness (mm)',
    'Insulin — 2-hour serum insulin (μU/mL)',
    'BMI — body mass index (kg/m²)',
    'DiabetesPedigreeFunction — family history score',
    'Age — age in years',
]
bullet_box(s, features, 0.6, 3.65, 12, 3.5, size=13)
slide_num(s, 10)
print('Slide 10 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — ML MODEL: PREPROCESSING
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Machine Learning: Data Preprocessing', 'Professional clinical data cleaning')
steps_pre = [
    ('Step 1', 'Zero Value Detection',
     'Medical columns (Glucose, BP, Insulin, BMI, SkinThickness)\ncannot be zero — replaced with NaN'),
    ('Step 2', 'Per-Class Median Imputation',
     'Missing values filled with class-specific median\n(diabetic median ≠ non-diabetic median)'),
    ('Step 3', 'Stratified Train/Test Split',
     '80% training (854 samples) / 20% test (214 samples)\nStratified to preserve class balance'),
    ('Step 4', 'Feature Scaling',
     'StandardScaler: zero mean, unit variance\nPrevents features with large ranges dominating'),
]
for i, (step, title, detail) in enumerate(steps_pre):
    y = 1.5 + i * 1.4
    rect(s, 0.4, y, 1.2, 1.2, MID_BLUE)
    txt(s, step, 0.45, y+0.35, 1.1, 0.5, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 1.8, y, 11.0, 1.2, RGBColor(0x1E,0x29,0x3B))
    txt(s, title, 1.95, y+0.05, 10.7, 0.45, size=14, bold=True, color=ACCENT_BLUE)
    txt(s, detail, 1.95, y+0.5, 10.7, 0.65, size=12, color=LIGHT_GRAY)
slide_num(s, 11)
print('Slide 11 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ML MODEL: RESULTS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Machine Learning: Model Performance', 'Gradient Boosting Classifier v2.3.1')
metric_box(s, 'Accuracy', '89.25%', 'test set', 0.4, 1.5, 2.8, 1.5, MID_BLUE, ACCENT_BLUE)
metric_box(s, 'ROC-AUC', '97.06%', 'excellent', 3.4, 1.5, 2.8, 1.5, RGBColor(0x05,0x46,0x2A), GREEN)
metric_box(s, 'Recall', '83.33%', 'catches diabetics', 6.4, 1.5, 2.8, 1.5, RGBColor(0x7C,0x3A,0xED), WHITE)
metric_box(s, 'F1 Score', '85.89%', 'balanced metric', 9.4, 1.5, 2.8, 1.5, ORANGE, DARK_BLUE)
txt(s, 'Cross-Validation (5-Fold):', 0.6, 3.3, 6, 0.4, size=14, bold=True, color=ACCENT_BLUE)
bullet_box(s, [
    'CV ROC-AUC Mean: 94.32%  (±1.53%)',
    'Overfitting Gap: 3.60%  (< 5% threshold ✅)',
    'Beats Logistic Regression by +8.27% AUC',
    'Beats Dummy Baseline by +44.25% AUC',
], 0.6, 3.75, 6, 2.5, size=13)
txt(s, 'Confusion Matrix (214 test samples):', 7.0, 3.3, 6, 0.4, size=14, bold=True, color=ACCENT_BLUE)
cm_data = [
    ('True Negatives', '122', 'Correctly identified non-diabetic', GREEN),
    ('False Positives', '8',  'False alarm (non-diabetic flagged)', ORANGE),
    ('False Negatives', '14', 'Missed diabetic cases', RED),
    ('True Positives', '70',  'Correctly identified diabetic', GREEN),
]
for i, (label, val, desc, color) in enumerate(cm_data):
    y = 3.75 + i * 0.65
    rect(s, 7.0, y, 0.6, 0.55, color)
    txt(s, val, 7.0, y+0.05, 0.6, 0.45, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, label + ': ' + desc, 7.75, y+0.1, 5.2, 0.4, size=12, color=LIGHT_GRAY)
slide_num(s, 12)
print('Slide 12 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — FEATURE IMPORTANCE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Feature Importance', 'Which health factors predict diabetes most?')
features_imp = [
    ('Insulin',                    47.3, ACCENT_BLUE),
    ('Glucose',                    20.0, GREEN),
    ('Age',                         9.1, RGBColor(0x7C,0x3A,0xED)),
    ('BMI',                         8.5, ORANGE),
    ('SkinThickness',               5.1, RGBColor(0x06,0xB6,0xD4)),
    ('DiabetesPedigreeFunction',    4.8, RGBColor(0xEC,0x48,0x99)),
    ('BloodPressure',               2.8, RGBColor(0x84,0xCC,0x16)),
    ('Pregnancies',                 2.4, RGBColor(0xF9,0x73,0x16)),
]
for i, (name, pct, color) in enumerate(features_imp):
    y = 1.5 + i * 0.67
    txt(s, name, 0.4, y+0.05, 3.5, 0.5, size=13, color=WHITE)
    bar_w = pct / 100 * 8.0
    rect(s, 4.0, y+0.05, bar_w, 0.45, color)
    txt(s, '%.1f%%' % pct, 4.1 + bar_w, y+0.1, 1.0, 0.35, size=13, bold=True, color=WHITE)
txt(s, 'Clinical insight: Insulin and Glucose are the strongest predictors — consistent with medical literature',
    0.4, 7.0, 12.5, 0.35, size=11, color=RGBColor(0x94,0xA3,0xB8), italic=True)
slide_num(s, 13)
print('Slide 13 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — RISK LEVELS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Risk Level Classification', '4-tier system with clinical recommendations')
levels = [
    ('🟢', 'LOW RISK', '0% – 30%', GREEN,
     'Very unlikely to have diabetes',
     'Maintain healthy lifestyle\nRegular checkups every 2–3 years'),
    ('🟡', 'MODERATE RISK', '30% – 50%', ORANGE,
     'Some risk factors present — monitor closely',
     'Consider dietary changes and exercise\nAnnual checkup recommended'),
    ('🟠', 'HIGH RISK', '50% – 70%', RGBColor(0xEA,0x58,0x0C),
     'Significant diabetes risk',
     'Consult healthcare provider soon\nSchedule appointment within 1 month'),
    ('🔴', 'VERY HIGH RISK', '70% – 100%', RED,
     'Strong likelihood of diabetes',
     'Immediate medical attention needed\nConsult doctor within 1 week'),
]
for i, (icon, level, prob, color, interp, action) in enumerate(levels):
    y = 1.5 + i * 1.4
    rect(s, 0.4, y, 2.5, 1.2, color)
    txt(s, icon + ' ' + level, 0.5, y+0.1, 2.3, 0.5, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, prob, 0.5, y+0.65, 2.3, 0.45, size=13, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 3.1, y, 4.5, 1.2, RGBColor(0x1E,0x29,0x3B))
    txt(s, interp, 3.2, y+0.3, 4.3, 0.6, size=12, color=LIGHT_GRAY)
    rect(s, 7.8, y, 5.1, 1.2, RGBColor(0x1E,0x29,0x3B))
    txt(s, action, 7.9, y+0.2, 4.9, 0.8, size=12, color=LIGHT_GRAY)
txt(s, 'Interpretation', 3.2, 1.2, 4.3, 0.3, size=12, bold=True, color=ACCENT_BLUE)
txt(s, 'Recommended Action', 7.9, 1.2, 4.9, 0.3, size=12, bold=True, color=ACCENT_BLUE)
slide_num(s, 14)
print('Slide 14 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — ROLES
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Role-Based Access Control', '6 roles — each with dedicated dashboard')
roles = [
    ('👤', 'Patient', ACCENT_BLUE,
     ['View health form', 'Submit prediction', 'View results', 'Book appointments', 'Pay online', 'View prescriptions']),
    ('🩺', 'Doctor', GREEN,
     ['View patient list', 'Review predictions', 'Write prescriptions', 'Order lab tests', 'Manage appointments', 'Clinical notes']),
    ('💉', 'Nurse', RGBColor(0x7C,0x3A,0xED),
     ['Register patients', 'Record vitals', 'Manage queue', 'View predictions', 'Clinical measurements', 'Patient notifications']),
    ('🔬', 'Lab Tech', ORANGE,
     ['View lab requests', 'Enter results', 'Validate results', 'Generate reports', 'Manage test types', 'Lab statistics']),
    ('💊', 'Pharmacist', RGBColor(0xEC,0x48,0x99),
     ['Review prescriptions', 'Check medication', 'Dispense drugs', 'Manage inventory', 'Expiry tracking', 'Dispensing history']),
    ('⚙️', 'Admin', RED,
     ['Manage all users', 'Approve payments', 'View audit logs', 'System reports', 'Model governance', 'Role management']),
]
for i, (icon, role, color, perms) in enumerate(roles):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.3
    y = 1.5 + row * 2.7
    rect(s, x, y, 3.9, 0.5, color)
    txt(s, icon + '  ' + role, x+0.1, y+0.07, 3.7, 0.38, size=15, bold=True, color=WHITE)
    bullet_box(s, perms, x+0.1, y+0.55, 3.7, 2.0, size=11, color=LIGHT_GRAY)
slide_num(s, 15)
print('Slide 15 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — NURSE WORKFLOW
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Nurse Workflow', 'Patient registration → vitals → queue management')
txt(s, 'When a new patient registers:', 0.6, 1.5, 12, 0.4, size=15, bold=True, color=ACCENT_BLUE)
bullet_box(s, [
    '1.  Patient registers (self or nurse-registered)',
    '2.  Nurse receives real-time notification: "New Patient Registered — PAT20260507XXXX"',
    '3.  Nurse clicks notification → Record Vitals page opens with patient pre-selected',
    '4.  Nurse records: Blood Pressure, BMI, Height/Weight, Skin Thickness, Age, Pregnancies',
    '5.  Patient removed from waiting queue automatically after vitals saved',
    '6.  Doctor notified: "Vitals recorded for [patient] — ready for consultation"',
], 0.6, 2.0, 12, 3.0, size=13)
txt(s, 'Queue Management:', 0.6, 5.2, 12, 0.4, size=15, bold=True, color=ACCENT_BLUE)
bullet_box(s, [
    'Waiting list shows all patients with human-readable wait time (e.g. "2 hr 15 min")',
    'Priority levels: Normal / Urgent / Emergency — color coded',
    'Auto-removed from queue when vitals are recorded',
], 0.6, 5.65, 12, 1.6, size=13)
slide_num(s, 16)
print('Slide 16 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — PREDICTION FLOW
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Prediction Flow', 'From health data to risk result')
flow = [
    ('1', 'Nurse Records\nVitals', 'BP, BMI, Age\nPregnancies\nSkin Thickness', MID_BLUE),
    ('2', 'Lab Tech\nEnters Glucose', 'Blood glucose\nresult from\nlab test', RGBColor(0x7C,0x3A,0xED)),
    ('3', 'Patient\nPays', 'Chapa / Cash\nInsurance\nDigital confirm', GREEN),
    ('4', 'Health Form\nAuto-Fills', 'All fields filled\nfrom DB\nPatient reviews', ACCENT_BLUE),
    ('5', 'ML Model\nPredicts', 'Gradient Boosting\n8 features\n< 1 second', ORANGE),
    ('6', 'Result\nDisplayed', 'Risk level\nProbability %\nRecommendation', RED),
]
for i, (num, title, detail, color) in enumerate(flow):
    x = 0.4 + i * 2.15
    rect(s, x, 1.5, 1.9, 0.6, color)
    txt(s, num, x+0.05, 1.55, 1.8, 0.5, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 5:
        txt(s, '→', x+1.9, 1.65, 0.25, 0.4, size=18, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    rect(s, x, 2.3, 1.9, 1.8, RGBColor(0x1E,0x29,0x3B))
    txt(s, title, x+0.1, 2.35, 1.7, 0.55, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, detail, x+0.1, 2.9, 1.7, 1.1, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
txt(s, 'After prediction: Doctor receives HIGH RISK alert → Reviews online → Writes digital prescription',
    0.4, 4.35, 12.5, 0.45, size=13, color=ACCENT_BLUE, italic=True)
txt(s, 'Key: Glucose ONLY comes from lab results — never manually entered. Ensures clinical accuracy.',
    0.4, 4.9, 12.5, 0.45, size=13, bold=True, color=ORANGE)
txt(s, 'Payment gate: Prediction only runs after payment is confirmed (Chapa auto-verify / Admin approves cash)',
    0.4, 5.45, 12.5, 0.45, size=13, color=LIGHT_GRAY)
slide_num(s, 17)
print('Slide 17 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — PAYMENT SYSTEM
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Payment Integration', 'Chapa · Cash · Insurance — all methods supported')
rect(s, 0.4, 1.5, 3.9, 5.5, MID_BLUE)
txt(s, '💳  Chapa (Online)', 0.55, 1.6, 3.6, 0.45, size=14, bold=True, color=WHITE)
bullet_box(s, [
    'TeleBirr',
    'CBE Birr',
    'M-Birr',
    'Visa / Mastercard',
    'Bank Transfer',
    '',
    'Auto-verified via API',
    'Webhook confirmation',
    'Instant prediction access',
], 0.55, 2.1, 3.6, 4.7, size=12, color=LIGHT_GRAY)
rect(s, 4.5, 1.5, 3.9, 5.5, RGBColor(0x05,0x46,0x2A))
txt(s, '💵  Cash', 4.65, 1.6, 3.6, 0.45, size=14, bold=True, color=WHITE)
bullet_box(s, [
    'Patient pays at cashier',
    'Gets reference number',
    'Admin approves in panel',
    'Patient clicks button:',
    '"Payment Approved —',
    ' Run My Prediction"',
    '',
    'Server confirms → runs',
    'prediction automatically',
], 4.65, 2.1, 3.6, 4.7, size=12, color=LIGHT_GRAY)
rect(s, 8.6, 1.5, 4.3, 5.5, RGBColor(0x1E,0x3A,0x8A))
txt(s, '🛡  Insurance', 8.75, 1.6, 4.0, 0.45, size=14, bold=True, color=WHITE)
bullet_box(s, [
    'Submit insurance claim',
    'Provider name + policy',
    'Claim reference number',
    'Admin reviews claim',
    'Approves in admin panel',
    '',
    'Patient notified via',
    'in-app notification',
    'with link to prediction',
], 8.75, 2.1, 4.0, 4.7, size=12, color=LIGHT_GRAY)
slide_num(s, 18)
print('Slide 18 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — NOTIFICATIONS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Real-Time Notification System', 'WebSocket push — no page refresh needed')
notifs = [
    ('🔴', 'HIGH RISK Alert', 'Sent to ALL doctors instantly when patient prediction is HIGH or VERY HIGH RISK', RED),
    ('🩺', 'New Patient Registered', 'Sent to all nurses with link to Record Vitals — patient ID and name pre-filled', MID_BLUE),
    ('🔬', 'Vitals Recorded', 'Sent to all doctors — patient ready for consultation', RGBColor(0x7C,0x3A,0xED)),
    ('📅', 'Appointment Confirmed', 'Sent to patient and doctor when appointment is booked', GREEN),
    ('✅', 'Payment Approved', 'Sent to patient when admin approves cash/insurance — includes link to run prediction', ORANGE),
    ('💊', 'Prescription Ready', 'Sent to patient when pharmacist dispenses medication', RGBColor(0xEC,0x48,0x99)),
]
for i, (icon, title, detail, color) in enumerate(notifs):
    y = 1.5 + i * 0.95
    rect(s, 0.4, y, 0.7, 0.8, color)
    txt(s, icon, 0.4, y+0.1, 0.7, 0.6, size=20, align=PP_ALIGN.CENTER)
    rect(s, 1.3, y, 11.6, 0.8, RGBColor(0x1E,0x29,0x3B))
    txt(s, title, 1.45, y+0.05, 4.0, 0.38, size=13, bold=True, color=WHITE)
    txt(s, detail, 5.6, y+0.1, 7.1, 0.6, size=12, color=LIGHT_GRAY)
slide_num(s, 19)
print('Slide 19 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — SECURITY
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Security & HIPAA Compliance', 'Healthcare-grade data protection')
sec = [
    ('🔑', 'JWT Authentication', '30-day tokens · Blacklist on logout · Role-based access · Token refresh'),
    ('🔒', 'Password Security', 'bcrypt hashing · Minimum 8 chars · Uppercase + number required'),
    ('🛡', 'Field Encryption', 'Fernet AES-128 encryption for PHI fields (medical history, allergies, medications)'),
    ('📋', 'Audit Logging', 'Every action logged: who, what, when, IP address · Immutable audit trail'),
    ('⚡', 'Rate Limiting', 'Login: 10/min · Register: 5/min · Prediction: 5/hour · Global: 5000/hour'),
    ('🌐', 'Security Headers', 'CSP · HSTS · X-Frame-Options · X-Content-Type-Options · Referrer-Policy'),
    ('🔐', 'Admin Protection', 'Admin signup disabled by default · Email lock · Separate admin token'),
    ('📧', 'OTP Verification', '6-digit OTP for email verification · 15-minute expiry · DB-stored'),
]
for i, (icon, title, detail) in enumerate(sec):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.5 + row * 1.35
    rect(s, x, y, 6.1, 1.2, RGBColor(0x1E,0x29,0x3B))
    txt(s, icon + '  ' + title, x+0.15, y+0.08, 5.8, 0.42, size=13, bold=True, color=ACCENT_BLUE)
    txt(s, detail, x+0.15, y+0.52, 5.8, 0.6, size=11, color=LIGHT_GRAY)
slide_num(s, 20)
print('Slide 20 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — TESTING OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Testing Strategy', '175+ automated tests across all layers')
metric_box(s, 'Total Tests', '175+', 'automated', 0.4, 1.5, 2.8, 1.4)
metric_box(s, 'ML Tests', '24/24', 'all passed ✅', 3.4, 1.5, 2.8, 1.4, RGBColor(0x05,0x46,0x2A), GREEN)
metric_box(s, 'System Tests', '20/22', 'passed ✅', 6.4, 1.5, 2.8, 1.4, MID_BLUE, ACCENT_BLUE)
metric_box(s, 'UAT Scenarios', '12/12', 'manual pass ✅', 9.4, 1.5, 2.8, 1.4, RGBColor(0x7C,0x3A,0xED), WHITE)
txt(s, 'Test Files:', 0.6, 3.2, 12, 0.4, size=14, bold=True, color=ACCENT_BLUE)
test_files = [
    'test_ml_model.py     — Data preprocessing, model training, prediction, evaluation (24 tests)',
    'test_comprehensive.py — Full system integration, performance, security headers (22 tests)',
    'test_routes.py       — All API endpoints: auth, patient, doctor, nurse, lab, pharmacy, admin (60+ tests)',
    'test_models.py       — Database models and route handlers (30+ tests)',
    'test_integration.py  — Cross-role workflows: patient → nurse → doctor → lab (10 tests)',
    'test_services.py     — Business logic services (15 tests)',
    'test_environment.py  — Environment configuration and dependencies (8 tests)',
    'test_db_config.py    — Database configuration and connectivity (6 tests)',
]
bullet_box(s, test_files, 0.6, 3.65, 12.5, 3.5, size=12)
slide_num(s, 21)
print('Slide 21 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — ML TEST RESULTS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'ML Model Tests — 24/24 PASSED ✅', 'All tests pass on every run')
test_results = [
    ('TestDataPreprocessing', ['test_load_data', 'test_data_no_missing_values', 'test_data_correct_types', 'test_feature_ranges'], GREEN),
    ('TestModelTraining', ['test_train_logistic_regression', 'test_train_random_forest', 'test_model_accuracy (≥75%)', 'test_model_roc_auc (≥0.85)'], ACCENT_BLUE),
    ('TestModelSaveLoad', ['test_save_model', 'test_load_model', 'test_save_scaler', 'test_load_scaler', 'test_save_features'], MID_BLUE),
    ('TestPrediction', ['test_prediction_output_type', 'test_prediction_probability_sum', 'test_risk_level_mapping', 'test_risk_color_mapping'], RGBColor(0x7C,0x3A,0xED)),
    ('TestModelEvaluation', ['test_confusion_matrix', 'test_precision_recall', 'test_f1_score'], ORANGE),
    ('TestMLServiceIntegration', ['test_ml_service_initialization', 'test_prediction_result_structure'], GREEN),
    ('TestFeatureImportance', ['test_feature_importance_random_forest', 'test_top_features'], ACCENT_BLUE),
]
y = 1.5
for (cls, tests, color) in test_results:
    rect(s, 0.4, y, 3.5, 0.38, color)
    txt(s, cls, 0.5, y+0.04, 3.3, 0.3, size=11, bold=True, color=WHITE)
    for j, t in enumerate(tests):
        x = 4.1 + j * 2.25
        if x < 12.8:
            rect(s, x, y, 2.1, 0.38, RGBColor(0x05,0x46,0x2A))
            txt(s, '✅ ' + t, x+0.05, y+0.04, 2.0, 0.3, size=9, color=WHITE)
    y += 0.52
slide_num(s, 22)
print('Slide 22 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — UAT SCENARIOS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'User Acceptance Testing (UAT)', '12 end-to-end scenarios — all passed')
uat = [
    ('Patient self-registration', 'Register → verify email → login → dashboard'),
    ('Nurse registers patient', 'Nurse login → register → notification to doctors'),
    ('Vitals recording', 'Select patient → record BP/BMI/age → queue updated'),
    ('Lab test flow', 'Doctor orders → lab enters result → patient notified'),
    ('Chapa payment → prediction', 'Select service → pay → verify → ML result'),
    ('Cash payment → admin → prediction', 'Pay cash → admin approves → button → prediction'),
    ('High-risk alert', 'High-risk data → all doctors notified instantly'),
    ('Prescription flow', 'Doctor prescribes → pharmacist dispenses → patient notified'),
    ('Password reset', 'Forgot password → OTP email → reset → login'),
    ('Multi-role isolation', 'Patient cannot access doctor routes'),
    ('Session expiry', 'Token expires → redirect to login → re-login works'),
    ('Notification system', 'WebSocket push → bell updates without refresh'),
]
for i, (scenario, steps) in enumerate(uat):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.5 + row * 0.95
    rect(s, x, y, 0.55, 0.82, GREEN)
    txt(s, '✅', x+0.05, y+0.18, 0.45, 0.45, size=16, align=PP_ALIGN.CENTER)
    rect(s, x+0.65, y, 5.65, 0.82, RGBColor(0x1E,0x29,0x3B))
    txt(s, scenario, x+0.75, y+0.04, 5.4, 0.38, size=12, bold=True, color=WHITE)
    txt(s, steps, x+0.75, y+0.44, 5.4, 0.35, size=10, color=LIGHT_GRAY)
slide_num(s, 23)
print('Slide 23 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — IMPACT METRICS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Impact Metrics', 'Before vs After comparison')
headers = ['Metric', 'Before', 'After']
rows = [
    ('Time to diagnosis',         '3–5 days',        'Same day ✅'),
    ('Undetected high-risk',      '~50% missed',     'Flagged automatically ✅'),
    ('Paper forms per visit',     '5–8 forms',       '0 (fully digital) ✅'),
    ('Payment processing',        '15–30 min cash',  '< 2 min (mobile) ✅'),
    ('Medical errors',            'Common',          'Eliminated ✅'),
    ('Doctor notification',       'Never',           'Real-time alert ✅'),
    ('Audit trail',               'None',            '100% logged ✅'),
    ('Prescription errors',       'Frequent',        'Zero (digital) ✅'),
    ('Lab result availability',   'Next day',        'Instant (same session) ✅'),
    ('Patient history access',    'Paper search',    'One click ✅'),
]
col_widths = [4.5, 3.5, 4.3]
col_starts = [0.4, 5.1, 8.8]
for j, (h, w) in enumerate(zip(headers, col_widths)):
    rect(s, col_starts[j], 1.5, w, 0.45, MID_BLUE)
    txt(s, h, col_starts[j]+0.1, 1.55, w-0.2, 0.35, size=13, bold=True, color=WHITE)
for i, row in enumerate(rows):
    y = 2.05 + i * 0.5
    bg_c = RGBColor(0x1E,0x29,0x3B) if i % 2 == 0 else RGBColor(0x0F,0x17,0x2A)
    for j, (val, w) in enumerate(zip(row, col_widths)):
        rect(s, col_starts[j], y, w, 0.48, bg_c)
        color = GREEN if '✅' in str(val) else (RED if 'missed' in str(val).lower() or 'common' in str(val).lower() or 'none' in str(val).lower() else WHITE)
        txt(s, val, col_starts[j]+0.1, y+0.07, w-0.2, 0.35, size=11, color=color)
slide_num(s, 24)
print('Slide 24 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — DEPLOYMENT
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Deployment & Production', 'Cloud-ready, scalable architecture')
rect(s, 0.4, 1.5, 5.9, 5.5, MID_BLUE)
txt(s, '☁️  Cloud Deployment (Render)', 0.55, 1.6, 5.6, 0.45, size=14, bold=True, color=WHITE)
bullet_box(s, [
    'Platform: Render.com (free tier)',
    'Database: PostgreSQL (production)',
    'SSL: Automatic HTTPS',
    'WSGI: Gunicorn server',
    'WebSocket: Flask-SocketIO',
    'Environment: .env variables',
    'Auto-deploy from GitHub',
    'Health check endpoint: /health',
], 0.55, 2.1, 5.6, 4.7, size=12, color=LIGHT_GRAY)
rect(s, 6.7, 1.5, 6.2, 5.5, RGBColor(0x1E,0x29,0x3B))
txt(s, '🐳  Docker Support', 6.85, 1.6, 5.9, 0.45, size=14, bold=True, color=WHITE)
bullet_box(s, [
    'Dockerfile included',
    'docker-compose.yml for local dev',
    'Environment variable injection',
    'Volume mounts for database',
    '',
    '📋  CI/CD Pipeline',
    'GitHub Actions workflow',
    'Automated test runs on push',
    'Lint and security checks',
    'Auto-deploy on main branch',
], 6.85, 2.1, 5.9, 4.7, size=12, color=LIGHT_GRAY)
slide_num(s, 25)
print('Slide 25 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — CHALLENGES & SOLUTIONS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Challenges & How We Solved Them', 'Real engineering problems encountered')
challenges = [
    ('Chapa Payment Verification',
     'Test sandbox does not confirm transactions via API',
     'Detect test key (CHASECK_TEST-) → skip API call → mark completed directly in DB'),
    ('Session Loss During Payment',
     'Chapa redirects back without valid JWT token',
     'Public verify endpoint + token refresh on payment_success.html load'),
    ('Cash Payment → Prediction Lock',
     'Patient pays cash but prediction stays locked after admin approves',
     'Server-side check on page load + "Check Payment" button + ?paid=1 URL param'),
    ('ML Model Overfitting',
     'Train accuracy 98.83% vs test 88.79% = 10% gap',
     'Increased regularization: max_depth 4→3, min_samples_leaf 4→10, added max_features=sqrt'),
    ('sklearn Version Mismatch',
     'Model saved with old sklearn, fails to load with 1.8.0',
     'Retrained and saved model with current sklearn version on deployment'),
    ('Queue Never Cleared',
     'Patients stayed in waiting list forever after vitals recorded',
     'record_vitals() now marks queue entry as completed after saving'),
]
for i, (title, problem, solution) in enumerate(challenges):
    y = 1.5 + i * 0.97
    rect(s, 0.4, y, 3.5, 0.85, MID_BLUE)
    txt(s, title, 0.5, y+0.05, 3.3, 0.75, size=11, bold=True, color=WHITE)
    rect(s, 4.1, y, 4.0, 0.85, RGBColor(0x7F,0x1D,0x1D))
    txt(s, '❌ ' + problem, 4.2, y+0.1, 3.8, 0.65, size=10, color=LIGHT_GRAY)
    rect(s, 8.3, y, 4.6, 0.85, RGBColor(0x05,0x46,0x2A))
    txt(s, '✅ ' + solution, 8.4, y+0.1, 4.4, 0.65, size=10, color=LIGHT_GRAY)
txt(s, 'Challenge', 0.5, 1.2, 3.3, 0.3, size=12, bold=True, color=ACCENT_BLUE)
txt(s, 'Problem', 4.2, 1.2, 3.8, 0.3, size=12, bold=True, color=RED)
txt(s, 'Solution', 8.4, 1.2, 4.4, 0.3, size=12, bold=True, color=GREEN)
slide_num(s, 26)
print('Slide 26 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Future Work & Improvements', 'Roadmap for next version')
future = [
    ('📊', 'Larger Dataset', 'Integrate CDC BRFSS or NHANES dataset (500,000+ samples)\nImprove model generalization across diverse populations'),
    ('🤖', 'Advanced ML', 'Experiment with XGBoost, LightGBM, Neural Networks\nEnsemble methods for higher recall on diabetic cases'),
    ('📱', 'Mobile App', 'React Native mobile app for patients\nOffline prediction capability for rural areas'),
    ('🔗', 'HL7 / FHIR', 'Healthcare interoperability standards\nIntegration with national health information systems'),
    ('📡', 'IoT Integration', 'Connect with glucometers and blood pressure monitors\nAuto-fill vitals from wearable devices'),
    ('🌍', 'Multi-Language', 'Amharic language support for Ethiopian patients\nLocalization for other African healthcare systems'),
]
for i, (icon, title, detail) in enumerate(future):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.5 + row * 1.9
    rect(s, x, y, 6.1, 1.7, RGBColor(0x1E,0x29,0x3B))
    txt(s, icon + '  ' + title, x+0.15, y+0.1, 5.8, 0.45, size=14, bold=True, color=ACCENT_BLUE)
    txt(s, detail, x+0.15, y+0.6, 5.8, 1.0, size=12, color=LIGHT_GRAY)
slide_num(s, 27)
print('Slide 27 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — KEY ACHIEVEMENTS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Key Technical Achievements', 'What makes this project stand out')
achievements = [
    ('89.25%', 'ML Accuracy', 'Gradient Boosting with 97.06% ROC-AUC'),
    ('6 Roles', 'RBAC System', 'Complete role-based access control'),
    ('Real-Time', 'Notifications', 'WebSocket push without page refresh'),
    ('3 Methods', 'Payments', 'Chapa + Cash + Insurance integrated'),
    ('175+', 'Tests', 'Automated test suite, 24/24 ML tests pass'),
    ('HIPAA', 'Compliant', 'Encryption, audit logs, rate limiting'),
    ('Same Day', 'Diagnosis', 'From registration to result in one session'),
    ('Zero Paper', 'Workflow', 'Fully digital clinical workflow'),
]
for i, (value, label, detail) in enumerate(achievements):
    col = i % 4
    row = i // 4
    x = 0.4 + col * 3.2
    y = 1.5 + row * 2.7
    rect(s, x, y, 2.9, 2.4, MID_BLUE)
    txt(s, value, x+0.1, y+0.15, 2.7, 0.75, size=28, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    txt(s, label, x+0.1, y+0.9, 2.7, 0.45, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, detail, x+0.1, y+1.4, 2.7, 0.85, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
slide_num(s, 28)
print('Slide 28 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 29 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
header_bar(s, 'Conclusion', 'AI can transform diabetes care in developing countries')
txt(s, 'What we built:', 0.6, 1.5, 12, 0.4, size=16, bold=True, color=ACCENT_BLUE)
bullet_box(s, [
    '✅  A production-ready AI-powered healthcare platform',
    '✅  ML model with 89.25% accuracy and 97.06% ROC-AUC for diabetes risk prediction',
    '✅  Complete role-based clinical workflow (6 roles, 175+ automated tests)',
    '✅  Real-time notifications, digital payments, HIPAA-compliant security',
    '✅  Same-day diagnosis — from registration to prediction result in one session',
], 0.6, 1.95, 12, 2.2, size=14)
txt(s, 'Problems solved:', 0.6, 4.3, 12, 0.4, size=16, bold=True, color=ACCENT_BLUE)
bullet_box(s, [
    '✅  Late diabetes diagnosis → Early AI-powered screening at every visit',
    '✅  Paper-based fragmented workflow → Fully digital, connected system',
    '✅  No decision support → Real-time ML risk scores for every patient',
    '✅  Disconnected departments → Real-time notifications between all roles',
    '✅  No digital payment → Chapa + Cash + Insurance with instant confirmation',
], 0.6, 4.75, 12, 2.2, size=14)
slide_num(s, 29)
print('Slide 29 done')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 30 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 7.5, DARK_BLUE)
rect(s, 0, 3.1, 13.33, 0.06, ACCENT_BLUE)
txt(s, '🙏', 5.9, 0.4, 1.5, 1.2, size=54, align=PP_ALIGN.CENTER)
txt(s, 'Thank You!', 0.5, 1.6, 12.3, 0.9, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, 'Questions & Discussion', 0.5, 2.55, 12.3, 0.5, size=18, color=ACCENT_BLUE, align=PP_ALIGN.CENTER, italic=True)
rect(s, 2.0, 3.4, 9.33, 2.5, MID_BLUE)
txt(s, '📌  Project Summary', 2.2, 3.5, 9.0, 0.45, size=15, bold=True, color=WHITE)
bullet_box(s, [
    '🤖  ML Model: Gradient Boosting · 89.25% Accuracy · 97.06% ROC-AUC',
    '👥  6 Roles: Patient · Doctor · Nurse · Lab · Pharmacist · Admin',
    '💳  Payments: Chapa · Cash · Insurance',
    '🔒  Security: JWT · bcrypt · Fernet encryption · HIPAA audit logs',
    '✅  Tests: 175+ automated · 24/24 ML tests pass',
], 2.2, 4.0, 9.0, 1.8, size=13, color=LIGHT_GRAY)
txt(s, 'GitHub: github.com/tasheayansa6/diabetes-prediction-system',
    0.5, 6.2, 12.3, 0.4, size=13, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
txt(s, 'Built with Flask · Python · scikit-learn · SQLite/PostgreSQL · Chapa API',
    0.5, 6.7, 12.3, 0.4, size=11, color=RGBColor(0x64,0x74,0x8B), align=PP_ALIGN.CENTER)
slide_num(s, 30)
print('Slide 30 done')

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
out = 'docs/Diabetes_Prediction_System_Presentation.pptx'
prs.save(out)
print()
print('=' * 60)
print('Presentation saved:', out)
print('Total slides:', len(prs.slides))
print('=' * 60)
